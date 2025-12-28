from __future__ import annotations

import hashlib
import tempfile
from pathlib import Path
from typing import List, Optional

from src.tools.document.parsers.types import DocumentContent, ImageContent, TableContent

try:
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
except Exception:
    Presentation = None  # type: ignore
    MSO_SHAPE_TYPE = None  # type: ignore


class PptxParser:
    """PPTX 演示文稿解析器。

    支持提取：
    - 幻灯片文本
    - 表格内容
    - 嵌入图片
    - 备注内容
    """

    def __init__(self, output_dir: Optional[Path] = None) -> None:
        """初始化解析器。

        Args:
            output_dir: 图片输出目录，默认使用临时目录
        """
        self._output_dir = output_dir

    def parse(self, path: Path) -> DocumentContent:
        """解析 PPTX 演示文稿。

        Args:
            path: PPTX 文件路径

        Returns:
            结构化的文档内容
        """
        if Presentation is None:
            return DocumentContent(
                text="",
                metadata={
                    "path": str(path),
                    "format": "pptx",
                    "warning": "python-pptx not installed.",
                },
            )

        try:
            prs = Presentation(str(path))

            slide_texts: List[str] = []
            all_images: List[ImageContent] = []
            all_tables: List[TableContent] = []

            for slide_idx, slide in enumerate(prs.slides, start=1):
                # 提取幻灯片文本
                texts = self._extract_slide_texts(slide)

                # 提取幻灯片备注
                notes = self._extract_notes(slide)
                if notes:
                    texts.append(f"[备注] {notes}")

                # 提取表格
                tables = self._extract_tables(slide, slide_idx)
                all_tables.extend(tables)

                # 提取图片
                images = self._extract_images(slide, path, slide_idx)
                all_images.extend(images)

                if texts:
                    slide_texts.append(f"[Slide {slide_idx}]\n" + "\n".join(texts))

            metadata = {
                "path": str(path),
                "format": "pptx",
                "parser": "python-pptx",
                "slides": len(prs.slides),
                "images": len(all_images),
                "tables": len(all_tables),
            }

            return DocumentContent(
                text="\n\n".join(slide_texts),
                images=all_images,
                tables=all_tables,
                metadata=metadata,
            )
        except Exception as exc:
            return DocumentContent(
                text="",
                metadata={
                    "path": str(path),
                    "format": "pptx",
                    "warning": f"PPTX parsing failed: {exc}",
                },
            )

    @staticmethod
    def _extract_slide_texts(slide) -> List[str]:
        """提取幻灯片中的文本。"""
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    texts.append(text)
        return texts

    @staticmethod
    def _extract_notes(slide) -> str:
        """提取幻灯片备注。"""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                if notes_frame and notes_frame.text:
                    return notes_frame.text.strip()
        except Exception:
            pass
        return ""

    @staticmethod
    def _extract_tables(slide, slide_idx: int) -> List[TableContent]:
        """提取幻灯片中的表格。"""
        tables: List[TableContent] = []

        if MSO_SHAPE_TYPE is None:
            return tables

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.has_table:
                try:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    tables.append(
                        TableContent(
                            raw=rows,
                            description=f"Slide {slide_idx} 表格 {shape_idx}",
                        )
                    )
                except Exception:
                    continue
        return tables

    def _extract_images(
        self, slide, doc_path: Path, slide_idx: int
    ) -> List[ImageContent]:
        """提取幻灯片中的图片。"""
        images: List[ImageContent] = []

        if MSO_SHAPE_TYPE is None:
            return images

        # 确定图片输出目录（优先使用配置的目录，否则使用文档同级 images 目录）
        output_dir = self._output_dir
        if output_dir is None:
            # 默认保存到文档同级的 images 目录，而不是临时目录
            output_dir = doc_path.parent / "images" / doc_path.stem
        output_dir.mkdir(parents=True, exist_ok=True)

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            try:
                # 检查是否是图片形状
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob

                    # 生成唯一文件名
                    content_hash = hashlib.md5(image_bytes).hexdigest()[:8]
                    ext = self._get_image_extension(image.content_type)
                    filename = f"{doc_path.stem}_s{slide_idx}_{shape_idx}_{content_hash}{ext}"
                    image_path = output_dir / filename

                    # 保存图片
                    image_path.write_bytes(image_bytes)

                    images.append(
                        ImageContent(
                            path=str(image_path),
                            caption="",
                            description="",
                            image_type="slide_image",
                        )
                    )
            except Exception:
                # 单个图片提取失败不影响其他图片
                continue

        return images

    @staticmethod
    def _get_image_extension(content_type: str) -> str:
        """根据 MIME 类型获取文件扩展名。"""
        type_map = {
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/gif": ".gif",
            "image/bmp": ".bmp",
            "image/tiff": ".tiff",
            "image/webp": ".webp",
            "image/x-emf": ".emf",
            "image/x-wmf": ".wmf",
        }
        return type_map.get(content_type, ".png")
