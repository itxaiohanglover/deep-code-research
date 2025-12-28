from __future__ import annotations

from pathlib import Path
from typing import List

from src.tools.document.parsers.types import TableContent

try:
    import docx  # type: ignore
except Exception:
    docx = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None  # type: ignore


class TableExtractor:
    """表格提取器。

    支持从多种文档格式中提取表格：
    - DOCX (Word 文档)
    - PPTX (PowerPoint 演示文稿)
    """

    def extract_docx(self, path: Path) -> List[TableContent]:
        """从 DOCX 文档中提取表格。

        Args:
            path: DOCX 文件路径

        Returns:
            表格内容列表
        """
        if docx is None:
            return []

        try:
            document = docx.Document(str(path))
            tables: List[TableContent] = []

            for idx, table in enumerate(document.tables, start=1):
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)

                # 尝试识别表头
                header = rows[0] if rows else []
                description = self._generate_table_description(
                    rows, header, f"DOCX 表格 {idx}"
                )

                tables.append(TableContent(raw=rows, description=description))

            return tables
        except Exception:
            return []

    def extract_pptx(self, path: Path) -> List[TableContent]:
        """从 PPTX 演示文稿中提取表格。

        Args:
            path: PPTX 文件路径

        Returns:
            表格内容列表
        """
        if Presentation is None:
            return []

        try:
            prs = Presentation(str(path))
            tables: List[TableContent] = []
            table_count = 0

            for slide_idx, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if shape.has_table:
                        try:
                            table_count += 1
                            table = shape.table
                            rows = []

                            for row in table.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                rows.append(cells)

                            # 尝试识别表头
                            header = rows[0] if rows else []
                            description = self._generate_table_description(
                                rows,
                                header,
                                f"Slide {slide_idx} 表格 {table_count}",
                            )

                            tables.append(
                                TableContent(raw=rows, description=description)
                            )
                        except Exception:
                            continue

            return tables
        except Exception:
            return []

    @staticmethod
    def _generate_table_description(
        rows: List[List[str]], header: List[str], default: str
    ) -> str:
        """生成表格描述。

        Args:
            rows: 表格行数据
            header: 表头行
            default: 默认描述

        Returns:
            表格描述字符串
        """
        if not rows:
            return default

        row_count = len(rows)
        col_count = len(rows[0]) if rows else 0

        # 构建描述
        desc_parts = [default]
        desc_parts.append(f"({row_count} 行 × {col_count} 列)")

        # 如果有明确的表头，添加列名信息
        if header and any(header):
            non_empty_headers = [h for h in header if h]
            if non_empty_headers:
                desc_parts.append(f"列: {', '.join(non_empty_headers[:5])}")
                if len(non_empty_headers) > 5:
                    desc_parts.append("...")

        return " ".join(desc_parts)

    def to_markdown(self, table: TableContent) -> str:
        """将表格转换为 Markdown 格式。

        Args:
            table: 表格内容

        Returns:
            Markdown 格式的表格字符串
        """
        if not table.raw:
            return ""

        rows = table.raw
        lines = []

        # 表头
        if rows:
            header = rows[0]
            lines.append("| " + " | ".join(str(cell) for cell in header) + " |")
            lines.append("| " + " | ".join("---" for _ in header) + " |")

            # 数据行
            for row in rows[1:]:
                lines.append("| " + " | ".join(str(cell) for cell in row) + " |")

        return "\n".join(lines)

    def to_csv(self, table: TableContent, delimiter: str = ",") -> str:
        """将表格转换为 CSV 格式。

        Args:
            table: 表格内容
            delimiter: 分隔符，默认逗号

        Returns:
            CSV 格式的字符串
        """
        if not table.raw:
            return ""

        lines = []
        for row in table.raw:
            # 处理包含分隔符或引号的单元格
            cells = []
            for cell in row:
                cell_str = str(cell)
                if delimiter in cell_str or '"' in cell_str or "\n" in cell_str:
                    cell_str = '"' + cell_str.replace('"', '""') + '"'
                cells.append(cell_str)
            lines.append(delimiter.join(cells))

        return "\n".join(lines)

